# This dependencies have to write for pptx module
import collections.abc
from pptx import Presentation

from src.slide import Slide
from src.convertors.presentation import PresentationConvertor


class PPTX:
    def __init__(self, filename: str):
        self.convertor = PresentationConvertor(Presentation())
        self.filename = filename
        self.slides = []
        self.page = 0

    def save(self):
        for slide in self.slides:
            self.convertor.add_slide(slide)
        self.convertor.save(self.filename)

    def add_slide(self, slide: Slide):
        self.slides.append(slide)
        self.page += 1
        return

    def page_num(self) -> int:
        return self.page
