from src.slide import Slide, StartPoint, Size
from src.components import TextBox, ListText, RecText
from src.placeholders import (
    AbstractPlaceHolder,
    PlaceHolderType,
    TitlePlaceHolder,
    ContentPlaceHolder,
    ListContentPlaceHolder,
)


from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE


PPTX_MAX_LEVEL = 8


class PlaceHolderPiceOfName:
    TITLE = "Title"
    CONTENT_PLACEHOLDER = "Content Placeholder"


def add_child_to_text_frame_rec(text_frame, parent: RecText, i: int):
    if parent is None:
        return
    if len(parent.children()) == 0:
        return

    # pptx only support 8 level list
    if i >= PPTX_MAX_LEVEL:
        i = PPTX_MAX_LEVEL - 1

    for child in parent.children():
        paragraph = text_frame.add_paragraph()
        paragraph.text = child.str()
        paragraph.font.bold = child.bold()
        paragraph.font.size = Pt(child.size())
        paragraph.level = i
        add_child_to_text_frame_rec(text_frame, child, i + 1)


def add_roots_to_text_frame(text_frame, roots: [RecText]):
    for root in roots:
        paragraph = text_frame.add_paragraph()
        paragraph.text = root.str()
        paragraph.font.bold = root.bold()
        paragraph.font.size = Pt(root.size())
        paragraph.level = 0
        add_child_to_text_frame_rec(text_frame, root, 1)


class PlaceHolderConvertor:
    def __init__(self, placeholders):
        self.placeholders = placeholders
        return

    def convert(self, placeholders: [AbstractPlaceHolder]):
        for placeholder in placeholders:
            if placeholder.type == PlaceHolderType.TITLE:
                self.__case_title(placeholder)

            if placeholder.type == PlaceHolderType.CONTENT:
                self.__case_content(placeholder)

            if placeholder.type == PlaceHolderType.LIST_CONTENT:
                self.__case_list_content(placeholder)

        return

    def __case_any(self, placeholder: AbstractPlaceHolder, pptx_placeholder_name, func):
        for pptx_placeholder in self.placeholders:
            if pptx_placeholder_name in pptx_placeholder.name:
                func(pptx_placeholder, placeholder)
                return

    def __case_title(self, placeholder: TitlePlaceHolder):
        def f(pptx_placeholder, placeholder):
            pptx_placeholder.text = placeholder.value
            return

        self.__case_any(placeholder, PlaceHolderPiceOfName.TITLE, f)

    def __case_content(self, placeholder: ContentPlaceHolder):
        def f(pptx_placeholder, placeholder):
            pptx_placeholder.text = placeholder.value[0]
            for i in range(1, len(placeholder.value)):
                para = pptx_placeholder.text_frame.add_paragraph()
                para.text = placeholder.value[i]
                return

        self.__case_any(placeholder, PlaceHolderPiceOfName.CONTENT_PLACEHOLDER, f)

    def __case_list_content(self, placeholder: ListContentPlaceHolder):
        def f(pptx_placeholder, placeholder):
            list_text: ListText = placeholder.value
            parents = list_text.lists()
            # set first parent to placeholder top text
            parent = parents[0]
            pptx_placeholder.text = parents[0].str()
            text_frame = pptx_placeholder.text_frame
            add_child_to_text_frame_rec(text_frame, parent, 1)

            # add rest of parents to text_frame
            add_roots_to_text_frame(text_frame, parents[1:])
            return

        self.__case_any(placeholder, PlaceHolderPiceOfName.CONTENT_PLACEHOLDER, f)


class SlideConvertor:
    def __init__(self, pptx_slide_api):
        self.pptx_slide_api = pptx_slide_api
        return

    def convert(self, slide: Slide):
        if slide is None:
            return
        self.__convert_placeholder(slide.placeholders)
        self.__convert_list_text(slide.list_texts)
        self.__convert_textbox(slide.textboxs)
        return

    def __convert_placeholder(self, placeholders: [AbstractPlaceHolder]):
        convertor = PlaceHolderConvertor(self.pptx_slide_api.placeholders)
        convertor.convert(placeholders)

    def __convert_list_text(self, list_texts: [ListText]):

        for list_text in list_texts:
            converted_box = self.__add_textbox(list_text.start_point, list_text.size)
            text_frame = converted_box.text_frame
            text_frame.text = ""

            if list_text is None or list_text.value is None:
                continue

            for top in list_text.value.lists():
                paragraph = text_frame.add_paragraph()
                paragraph.text = top.str()
                paragraph.font.bold = top.bold()
                paragraph.font.size = Pt(top.size())
                paragraph.level = 0
                add_child_to_text_frame_rec(text_frame, top, 1)
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        return

    def __convert_textbox(self, textboxes: [TextBox]):
        for textbox in textboxes:
            converted_box = self.__add_textbox(textbox.start_point, textbox.size)
            text_frame = converted_box.text_frame
            text_frame.text = ""

            if textbox is None or textbox.value is None:
                continue

            for i, text in enumerate(textbox.value.texts):
                paragraph = text_frame.add_paragraph()
                paragraph.text = text.str()
                paragraph.font.bold = text.bold
                paragraph.font.size = Pt(text.size())
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        return

    def __add_textbox(self, start_point: StartPoint, size: Size):
        return self.pptx_slide_api.shapes.add_textbox(
            Pt(start_point.left),
            Pt(start_point.top),
            Pt(size.width),
            Pt(size.height),
        )
