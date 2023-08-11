from src.slide import Slide, StartPoint, Size
from src.components import TextBox, ListText, RecText
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE


class SlideConvertor:
    def __init__(self, pptx_slide_api):
        self.pptx_slide_api = pptx_slide_api
        return

    def convert(self, slide: Slide):
        self.__convert_list_text(slide.list_texts)
        self.__convert_textbox(slide.textboxs)
        return

    def __convert_list_text(self, list_texts: [ListText]):
        def children(text_frame, parent: RecText, i: int):
            if len(parent.children()) == 0:
                return

            for child in parent.children():
                paragraph = text_frame.add_paragraph()
                paragraph.text = child.str()
                paragraph.font.bold = child.bold()
                paragraph.font.size = Pt(child.size())
                paragraph.level = i
                children(text_frame, child, i + 1)

        for list_text in list_texts:
            converted_box = self.__add_textbox(list_text.start_point, list_text.size)
            text_frame = converted_box.text_frame
            text_frame.text = ""

            if list_text is None or list_text.value is None:
                continue

            for top in list_text.value.lists():
                paragraph = text_frame.add_paragraph()
                paragraph.text = top.str()
                paragraph.font.bold = top.bold()
                paragraph.font.size = Pt(top.size())
                paragraph.level = 0
                children(text_frame, top, 1)
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        return

    def __convert_textbox(self, textboxes: [TextBox]):
        for textbox in textboxes:
            converted_box = self.__add_textbox(textbox.start_point, textbox.size)
            text_frame = converted_box.text_frame
            text_frame.text = ""

            if textbox is None or textbox.value is None:
                continue

            for i, text in enumerate(textbox.value.texts):
                paragraph = text_frame.add_paragraph()
                paragraph.text = text.str()
                paragraph.font.bold = text.bold
                paragraph.font.size = Pt(text.size())
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        return

    def __add_textbox(self, start_point: StartPoint, size: Size):
        return self.pptx_slide_api.shapes.add_textbox(
            Pt(start_point.left),
            Pt(start_point.top),
            Pt(size.width),
            Pt(size.height),
        )
