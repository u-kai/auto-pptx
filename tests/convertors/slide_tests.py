from src.convertors.slide import SlideConvertor
from src.components import TextBox, Text, Font, ListText
from pptx.util import Pt
from src.slide import Slide, StartPoint, Size
from src.placeholders import TitlePlaceHolder,ContentPlaceHolder,ListContentPlaceHolder
import unittest


class MockPPTXFont:
    def __init__(self):
        self.bold = False
        self.size = Pt(18)
        return


class MockPPTXParagraph:
    def __init__(self):
        self.font = MockPPTXFont()
        self.text = ""
        self.level = 0
        return


class MockPPTXTextFrame:
    def __init__(self):
        self.text = ""
        self.paragraphs = []
        return

    def add_paragraph(self):
        para = MockPPTXParagraph()
        self.paragraphs.append(para)
        return para


class MockPPTXTextBox:
    def __init__(
        self,
    ):
        self.text_frame = MockPPTXTextFrame()
        return


class MockPPTXShapesApi:
    def __init__(self):
        self.textboxs = []

    def add_textbox(self, left, top, width, height) -> MockPPTXTextBox:
        textbox = MockPPTXTextBox()
        self.textboxs.append(
            {
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "textbox": textbox,
            }
        )
        return textbox


class MockSubTitlePlaceHolder():
    def __init__(self):
        self.name = "SubTitle 2"
        self.text = ""
        return


class MockContntPlaceHolder():
    def __init__(self,num:int):
        self.name = "Content Placeholder "+str(num)
        self.text = ""
        self.text_frame = MockPPTXTextFrame()
        return

class MockTitlePlaceHolder():
    def __init__(self):
        self.name = "Title 1"
        self.text = ""
        return


class MockPPTXSlideApi:
    def __init__(self):
        self.placeholders = []
        self.shapes = MockPPTXShapesApi()
        return


class TestSlideConvertorPlaceHolder(unittest.TestCase):
    def test_slideからlist_contentを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        mock.placeholders = [MockTitlePlaceHolder(),MockContntPlaceHolder(2)]
        # My types 
        list_content = ListContentPlaceHolder()
        root = Text("Root")
        root.change_size(28)
        root.to_bold()
        list_content.add(root)

        parent = Text("Parent1")
        parent.change_size(24)
        list_content.add_child_to(0, parent)

        list_content.top(0).child(0).add_child(Text("Child1"))
        slide = Slide.title_and_content()
        slide.add_placeholder(list_content)
        sut = SlideConvertor(mock)

        sut.convert(slide)

        text_frame:MockPPTXTextFrame = mock.placeholders[1].text_frame

        self.assertEqual(text_frame.paragraphs[0].text, "Root")
        self.assertEqual(text_frame.paragraphs[0].font.size, Pt(28))
        self.assertEqual(text_frame.paragraphs[0].font.bold, True)
        self.assertEqual(text_frame.paragraphs[0].level, 0)

        self.assertEqual(text_frame.paragraphs[1].text, "Parent1")
        self.assertEqual(text_frame.paragraphs[1].font.size, Pt(24))
        self.assertEqual(text_frame.paragraphs[1].font.bold, False)
        self.assertEqual(text_frame.paragraphs[1].level, 1)

        self.assertEqual(text_frame.paragraphs[2].text, "Child1")
        self.assertEqual(text_frame.paragraphs[2].font.size, Pt(18))
        self.assertEqual(text_frame.paragraphs[2].font.bold, False)
        self.assertEqual(text_frame.paragraphs[2].level, 2)

    def test_slideからcontent_placeholderを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        mock.placeholders = [MockTitlePlaceHolder(),MockContntPlaceHolder(2)]
        # My types
        contents = ContentPlaceHolder()
        contents.add("Hello World")
        contents.add("Good Bye")
        slide = Slide.title_and_content()
        slide.add_placeholder(contents)

        sut = SlideConvertor(mock)

        sut.convert(slide)
        
        self.assertEqual(mock.placeholders[1].text, "Hello World")
        text_frame:MockPPTXTextFrame = mock.placeholders[1].text_frame
        self.assertEqual(text_frame.paragraphs[0].text, "Good Bye")

    def test_slideからtitle_placeholderを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        mock.placeholders = [MockTitlePlaceHolder(),MockSubTitlePlaceHolder()]
        # My types
        title = TitlePlaceHolder("Hello World")
        slide = Slide.title_slide()
        slide.add_placeholder(title)

        sut = SlideConvertor(mock)

        sut.convert(slide)

        self.assertEqual(mock.placeholders[0].text, "Hello World")



class TestSlideConvertorListTexts(unittest.TestCase):
    def test_slideからtext_listを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Root")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        list = ListText(text)

        slide = Slide()
        slide.add_list_text(StartPoint(0, 0), Size(300, 300), list)

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], Pt(300))
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text, "Root"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold,
            True,
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].level, 0
        )

    def test_slideからtext_listを再起的に取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Root")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        list = ListText(text)
        text = Text("Root2")
        child = Text("Parent2")
        list.add_siblings(text)
        list.add_child_to(1, child)

        text = Text("Parent")
        list.add_child_to(0, text)

        text = Text("Child")
        list.top(0).child(0).add_child(text)

        text = Text("GrandChild")
        list.top(0).child(0).child(0).add_child(text)

        slide = Slide()
        slide.add_list_text(StartPoint(0, 0), Size(300, 300), list)

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], Pt(300))
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text, "Root"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold,
            True,
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].level, 0
        )

        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].text, "Parent"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.bold,
            False,
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].level, 1
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[2].text, "Child"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[2].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[2].font.bold,
            False,
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[2].level, 2
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[3].text,
            "GrandChild",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[3].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[3].font.bold,
            False,
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[3].level, 3
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[4].text, "Root2"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[4].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[4].font.bold,
            False,
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[4].level, 0
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[5].text,
            "Parent2",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[5].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[5].font.bold,
            False,
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[5].level, 1
        )



class TestSlideConvertorTextboxs(unittest.TestCase):
    def test_slideからtextboxを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Hello World")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        box = TextBox()
        box.add(text)

        slide = Slide()
        slide.add_textbox(
            StartPoint(0, 0),
            Size(300, 300),
            box,
        )

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], Pt(300))
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text,
            "Hello World",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold, True
        )

    def test_slideから複数textが格納されているtextboxを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Hello World")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()
        box = TextBox()
        box.add(text)

        text = Text("Good Bye")
        font = Font.meiryo()
        font.change_size(18)
        text.change_font(font)
        box.add(text)

        slide = Slide()
        slide.add_textbox(
            StartPoint(0, 1000),
            Size(300, 0),
            box,
        )

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], Pt(1000))
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], 0)
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text,
            "Hello World",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold, True
        )

        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].text, "Good Bye"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.bold, False
        )

    def test_slideからtextboxを複数取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Hello World")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        text2 = Text("Good Bye")
        font = Font.meiryo()
        font.change_size(18)
        text2.change_font(font)

        box = TextBox()
        box.add(text)
        box.add(text2)

        slide = Slide()
        slide.add_textbox(
            StartPoint(0, 0),
            Size(300, 300),
            box,
        )

        text = Text("Rust Good Language")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        text2 = Text("Python Good Language.But I love Rust more than Python")
        font = Font.meiryo()
        font.change_size(18)
        text2.change_font(font)

        box = TextBox()
        box.add(text)
        box.add(text2)

        slide.add_textbox(
            StartPoint(0, 1000),
            Size(100, 500),
            box,
        )

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], Pt(300))
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text,
            "Hello World",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold, True
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].text, "Good Bye"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.bold, False
        )

        self.assertEqual(mock.shapes.textboxs[1]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[1]["top"], Pt(1000))
        self.assertEqual(mock.shapes.textboxs[1]["width"], Pt(100))
        self.assertEqual(mock.shapes.textboxs[1]["height"], Pt(500))
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[0].text,
            "Rust Good Language",
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[0].font.bold, True
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[1].text,
            "Python Good Language.But I love Rust more than Python",
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[1].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[1].font.bold, False
        )


if __name__ == "__main__":
    unittest.main()
