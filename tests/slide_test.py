import unittest
from src.slide import Slide, StartPoint, Size
from src.convertor import SlideConvertor
from src.components import TextBox, Text, Font, ListText
from pptx.util import Pt


class TestSlide(unittest.TestCase):
    def test_slideにtext_boxを追加可能(self):

        sut = Slide()
        box = TextBox()
        box.add(Text("Hello World"))
        left = 0
        top = 0
        point = StartPoint(left, top)
        width = 300
        height = 300
        size = Size(width, height)

        sut.add_textbox(point, size, box)

        self.assertEqual(sut.textboxs[0].start_point, point)
        self.assertEqual(sut.textboxs[0].size, size)
        self.assertEqual(sut.textboxs[0].value, box)

    def test_slideにlist_textを追加可能(self):
        sut = Slide()
        list_text = ListText(Text("Root"))
        list_text.add_child_to(0, Text("Parent"))
        list_text.top(0).child(0).add_child(Text("Child"))
        left = 0
        top = 0
        point = StartPoint(left, top)
        width = 300
        height = 300
        size = Size(width, height)

        sut.add_list_text(point, size, list_text)

        self.assertEqual(sut.list_texts[0].start_point, point)
        self.assertEqual(sut.list_texts[0].size, size)
        self.assertEqual(sut.list_texts[0].value, list_text)


class MockPPTXFont:
    def __init__(self):
        self.bold = False
        self.size = Pt(18)
        return


class MockPPTXParagraph:
    def __init__(self):
        self.font = MockPPTXFont()
        self.text = ""
        self.level = 0
        return


class MockPPTXTextFrame:
    def __init__(self):
        self.paragraphs = []
        return

    def add_paragraph(self):
        para = MockPPTXParagraph()
        self.paragraphs.append(para)
        return para


class MockPPTXTextBox:
    def __init__(
        self,
    ):
        self.text_frame = MockPPTXTextFrame()
        return


class MockPPTXShapesApi:
    def __init__(self):
        self.textboxs = []

    def add_textbox(self, left, top, width, height) -> MockPPTXTextBox:
        textbox = MockPPTXTextBox()
        self.textboxs.append(
            {
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "textbox": textbox,
            }
        )
        return textbox


class MockPPTXSlideApi:
    def __init__(self):
        self.shapes = MockPPTXShapesApi()
        return


class TestSlideConvertor_textboxs(unittest.TestCase):
    def test_slideからtextboxを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Hello World")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        box = TextBox()
        box.add(text)

        slide = Slide()
        slide.add_textbox(
            StartPoint(0, 0),
            Size(300, 300),
            box,
        )

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], Pt(300))
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text,
            "Hello World",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold, True
        )

    def test_slideから複数textが格納されているtextboxを取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Hello World")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()
        box = TextBox()
        box.add(text)

        text = Text("Good Bye")
        font = Font.meiryo()
        font.change_size(18)
        text.change_font(font)
        box.add(text)

        slide = Slide()
        slide.add_textbox(
            StartPoint(0, 1000),
            Size(300, 0),
            box,
        )

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], Pt(1000))
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], 0)
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text,
            "Hello World",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold, True
        )

        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].text, "Good Bye"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.bold, False
        )

    def test_slideからtextboxを複数取得して変換可能(self):
        # 3rd party library mock
        mock = MockPPTXSlideApi()
        # My types
        text = Text("Hello World")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        text2 = Text("Good Bye")
        font = Font.meiryo()
        font.change_size(18)
        text2.change_font(font)

        box = TextBox()
        box.add(text)
        box.add(text2)

        slide = Slide()
        slide.add_textbox(
            StartPoint(0, 0),
            Size(300, 300),
            box,
        )

        text = Text("Rust Good Language")
        font = Font.meiryo_ui()
        font.change_size(28)
        text.change_font(font)
        text.to_bold()

        text2 = Text("Python Good Language.But I love Rust more than Python")
        font = Font.meiryo()
        font.change_size(18)
        text2.change_font(font)

        box = TextBox()
        box.add(text)
        box.add(text2)

        slide.add_textbox(
            StartPoint(0, 1000),
            Size(100, 500),
            box,
        )

        sut = SlideConvertor(mock)
        sut.convert(slide)

        self.assertEqual(mock.shapes.textboxs[0]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["top"], 0)
        self.assertEqual(mock.shapes.textboxs[0]["width"], Pt(300))
        self.assertEqual(mock.shapes.textboxs[0]["height"], Pt(300))
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].text,
            "Hello World",
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[0].font.bold, True
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].text, "Good Bye"
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[0]["textbox"].text_frame.paragraphs[1].font.bold, False
        )

        self.assertEqual(mock.shapes.textboxs[1]["left"], 0)
        self.assertEqual(mock.shapes.textboxs[1]["top"], Pt(1000))
        self.assertEqual(mock.shapes.textboxs[1]["width"], Pt(100))
        self.assertEqual(mock.shapes.textboxs[1]["height"], Pt(500))
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[0].text,
            "Rust Good Language",
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[0].font.size,
            Pt(28),
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[0].font.bold, True
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[1].text,
            "Python Good Language.But I love Rust more than Python",
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[1].font.size,
            Pt(18),
        )
        self.assertEqual(
            mock.shapes.textboxs[1]["textbox"].text_frame.paragraphs[1].font.bold, False
        )


if __name__ == "__main__":
    unittest.main()
